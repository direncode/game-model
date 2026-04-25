"""PPTX → text for the NATO Simulation corpus prep.

Works on .pptx (OOXML). Legacy .ppt must be converted first — see
``scripts/convert_ppt_to_pptx.py`` or just open and re-save via
PowerPoint/LibreOffice.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class Slide:
    number: int
    text: str


def extract_pptx_slides(path: str | Path) -> list[Slide]:
    """Return one ``Slide`` per slide, in order, with concatenated shape text."""
    from pptx import Presentation  # local import

    p = Presentation(str(path))
    out: list[Slide] = []
    for i, slide in enumerate(p.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            if text.strip():
                parts.append(text.strip())
        out.append(Slide(number=i, text="\n".join(parts)))
    return out


def pptx_as_single_document(path: str | Path) -> str:
    """Return all slides as a newline-joined string, with slide markers."""
    slides = extract_pptx_slides(path)
    return "\n\n".join(f"[Slide {s.number}]\n{s.text}" for s in slides)
